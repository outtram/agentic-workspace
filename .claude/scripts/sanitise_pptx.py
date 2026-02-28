#!/usr/bin/env python3
"""
Sanitise PPTX files by replacing sensitive client/company names.

Usage:
    python3 .claude/scripts/sanitise_pptx.py input.pptx [output.pptx]
    python3 .claude/scripts/sanitise_pptx.py input.pptx --rules rules.yml
    python3 .claude/scripts/sanitise_pptx.py input.pptx --dry-run

Default rules file: .claude/scripts/sanitise_rules.yml
If no output path given, saves as input_sanitised.pptx
"""

import argparse
import copy
import os
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu


DEFAULT_RULES_PATH = Path(__file__).parent / "sanitise_rules.yml"


def load_rules(rules_path: str | None = None) -> list[dict]:
    path = Path(rules_path) if rules_path else DEFAULT_RULES_PATH
    if not path.exists():
        print(f"Rules file not found: {path}", file=sys.stderr)
        sys.exit(1)
    with open(path) as f:
        data = yaml.safe_load(f)
    return data.get("replacements", [])


def apply_replacements(text: str, rules: list[dict]) -> str:
    for rule in rules:
        pattern = rule["pattern"]
        replacement = rule["replacement"]
        flags = re.IGNORECASE if rule.get("case_insensitive", True) else 0
        text = re.sub(pattern, replacement, text, flags=flags)
    return text


def sanitise_run(run, rules: list[dict]) -> int:
    if not run.text:
        return 0
    original = run.text
    sanitised = apply_replacements(original, rules)
    if sanitised != original:
        run.text = sanitised
        return 1
    return 0


def sanitise_paragraph(paragraph, rules: list[dict]) -> int:
    count = 0
    for run in paragraph.runs:
        count += sanitise_run(run, rules)
    return count


def sanitise_text_frame(text_frame, rules: list[dict]) -> int:
    count = 0
    for paragraph in text_frame.paragraphs:
        count += sanitise_paragraph(paragraph, rules)
    return count


def sanitise_table(table, rules: list[dict]) -> int:
    count = 0
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                count += sanitise_text_frame(cell.text_frame, rules)
    return count


def sanitise_shape(shape, rules: list[dict]) -> int:
    count = 0
    if shape.has_text_frame:
        count += sanitise_text_frame(shape.text_frame, rules)
    if shape.has_table:
        count += sanitise_table(shape.table, rules)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            count += sanitise_shape(child, rules)
    return count


def sanitise_notes(slide, rules: list[dict]) -> int:
    count = 0
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame
        if notes:
            count += sanitise_text_frame(notes, rules)
    return count


def sanitise_pptx(input_path: str, output_path: str, rules: list[dict], dry_run: bool = False) -> dict:
    prs = Presentation(input_path)
    stats = {"slides": 0, "replacements": 0, "slides_affected": []}

    for i, slide in enumerate(prs.slides, 1):
        slide_count = 0
        for shape in slide.shapes:
            slide_count += sanitise_shape(shape, rules)
        slide_count += sanitise_notes(slide, rules)

        if slide_count > 0:
            stats["slides_affected"].append(i)
            stats["replacements"] += slide_count
        stats["slides"] += 1

    # Also sanitise slide masters and layouts
    master_count = 0
    for master in prs.slide_masters:
        for shape in master.shapes:
            master_count += sanitise_shape(shape, rules)
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                master_count += sanitise_shape(shape, rules)
    stats["master_replacements"] = master_count
    stats["replacements"] += master_count

    if not dry_run:
        prs.save(output_path)
        stats["output"] = output_path

    return stats


def main():
    parser = argparse.ArgumentParser(description="Sanitise PPTX files")
    parser.add_argument("input", help="Input PPTX file path")
    parser.add_argument("output", nargs="?", help="Output PPTX file path")
    parser.add_argument("--rules", help="YAML rules file path")
    parser.add_argument("--dry-run", action="store_true", help="Count replacements without saving")
    args = parser.parse_args()

    if not args.output:
        stem = Path(args.input).stem
        parent = Path(args.input).parent
        args.output = str(parent / f"{stem}_sanitised.pptx")

    rules = load_rules(args.rules)
    stats = sanitise_pptx(args.input, args.output, rules, args.dry_run)

    print(f"Slides scanned: {stats['slides']}")
    print(f"Replacements made: {stats['replacements']} (inc. {stats.get('master_replacements', 0)} in masters/layouts)")
    print(f"Slides affected: {stats['slides_affected']}")
    if not args.dry_run:
        print(f"Saved to: {stats['output']}")
    else:
        print("(dry run — no file saved)")


if __name__ == "__main__":
    main()
