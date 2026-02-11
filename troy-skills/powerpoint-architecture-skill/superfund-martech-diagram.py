#!/usr/bin/env python3
"""
Super Fund MarTech Architecture Diagram
Using pptx-arch-diagrams skill v1.2 guidelines
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# v1.2 Color Scheme
COLORS = {
    'blue': RGBColor(0, 120, 212),      # API Gateway
    'teal': RGBColor(0, 128, 128),      # Core Services (v1.2)
    'orange': RGBColor(255, 140, 0),    # UI/Channels
    'green': RGBColor(16, 124, 16),     # Databases
    'purple': RGBColor(92, 45, 145),    # External/CMS
    'red': RGBColor(232, 17, 35),       # Security
    'gray': RGBColor(118, 118, 118),    # Infrastructure
}

# v1.2 Layout Constants - CONSISTENT ALIGNMENT
LEFT_MARGIN = 0.5
RIGHT_MARGIN = 0.5
CONTENT_WIDTH = 9.0
GAP = 0.08

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color, text_size=12, text_color=RGBColor(255, 255, 255)):
    """Add a shape with text and styling"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(text_size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = text_color

    return shape

def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(100, 100, 100)):
    """Add arrow connector"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def add_label(slide, text, x, y, width=2.5, height=0.25, size=11, bold=True, color=RGBColor(0, 0, 0)):
    """Add text label"""
    label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    label.text_frame.text = text
    p = label.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return label

# ===== SLIDE 1: TITLE SLIDE =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide1.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(2), Inches(CONTENT_WIDTH), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Super Fund MarTech Architecture"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

subtitle_box = slide1.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(3.5), Inches(CONTENT_WIDTH), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Multi-Channel Digital Ecosystem - Future State"
p2 = subtitle_frame.paragraphs[0]
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(100, 100, 100)

# ===== SLIDE 2: ARCHITECTURE OVERVIEW =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
add_label(slide2, "MarTech Architecture Overview", LEFT_MARGIN, 0.2, CONTENT_WIDTH, 0.4, 28, True)

# LAYER 1: Channel Layer (Orange) - Member touchpoints
y1 = 0.75
add_label(slide2, "Channel Layer", LEFT_MARGIN, y1, 2.5, 0.25, 11)
y1_comps = y1 + 0.3
comp_width = (CONTENT_WIDTH - (3 * GAP)) / 4

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN, y1_comps, comp_width, 0.45,
                   "Public Website", COLORS['orange'], 11)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + (comp_width + GAP), y1_comps, comp_width, 0.45,
                   "Member Portal", COLORS['orange'], 11)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + 2*(comp_width + GAP), y1_comps, comp_width, 0.45,
                   "iOS App", COLORS['orange'], 11)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + 3*(comp_width + GAP), y1_comps, comp_width, 0.45,
                   "Android App", COLORS['orange'], 11)

# Arrow
arrow_y1 = y1_comps + 0.55
add_arrow(slide2, 5, arrow_y1, 5, arrow_y1 + 0.12)

# LAYER 2: API Gateway (Blue)
y2 = arrow_y1 + 0.15
add_label(slide2, "API Layer", LEFT_MARGIN, y2, 2.5, 0.25, 11)
y2_comps = y2 + 0.3

add_shape_with_text(slide2, MSO_SHAPE.HEXAGON,
                   LEFT_MARGIN, y2_comps, CONTENT_WIDTH, 0.4,
                   "API Gateway", COLORS['blue'], 13)

# Arrow
arrow_y2 = y2_comps + 0.5
add_arrow(slide2, 5, arrow_y2, 5, arrow_y2 + 0.12)

# LAYER 3: MarTech Services (Teal - v1.2)
y3 = arrow_y2 + 0.15
add_label(slide2, "MarTech Services Layer", LEFT_MARGIN, y3, 2.5, 0.25, 11)
y3_comps = y3 + 0.3
comp_width3 = (CONTENT_WIDTH - (5 * GAP)) / 6

services = ["Member\nServices", "Content\nMgmt", "Campaign\nMgmt", "Marketing\nAutomation", "Analytics", "Personalization"]
for i, service in enumerate(services):
    add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(comp_width3 + GAP), y3_comps,
                       comp_width3, 0.4,
                       service, COLORS['teal'], 9)

# Arrow
arrow_y3 = y3_comps + 0.5
add_arrow(slide2, 5, arrow_y3, 5, arrow_y3 + 0.12)

# LAYER 4: Data Platform (Green)
y4 = arrow_y3 + 0.15
add_label(slide2, "Data Platform Layer", LEFT_MARGIN, y4, 2.5, 0.25, 11)
y4_comps = y4 + 0.3
comp_width4 = (CONTENT_WIDTH - (3 * GAP)) / 4

databases = ["Member DB", "Marketing DB", "Analytics\nPlatform", "CDP"]
for i, db in enumerate(databases):
    add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(comp_width4 + GAP), y4_comps,
                       comp_width4, 0.38,
                       db, COLORS['green'], 10)

# Footer
add_label(slide2, "v1.2 | Teal Services | Consistent Margins | Full-Width Alignment",
         LEFT_MARGIN, 5.45, CONTENT_WIDTH, 0.15, 8, False, RGBColor(150, 150, 150))

# ===== SLIDE 3: DETAILED VIEW WITH INTEGRATIONS =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

add_label(slide3, "MarTech Architecture - Integration Layer", LEFT_MARGIN, 0.2, CONTENT_WIDTH, 0.4, 28, True)

# Channels (smaller, top)
y_ch = 0.75
add_label(slide3, "Channels", LEFT_MARGIN, y_ch, 2, 0.2, 10)
y_ch_c = y_ch + 0.25
ch_width = (CONTENT_WIDTH - (3 * GAP)) / 4
channels = ["Web", "Portal", "iOS", "Android"]
for i, ch in enumerate(channels):
    add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(ch_width + GAP), y_ch_c,
                       ch_width, 0.35,
                       ch, COLORS['orange'], 10)

# API Gateway
y_api = y_ch_c + 0.45
add_arrow(slide3, 5, y_api - 0.05, 5, y_api + 0.05)
add_shape_with_text(slide3, MSO_SHAPE.HEXAGON,
                   LEFT_MARGIN, y_api + 0.1, CONTENT_WIDTH, 0.35,
                   "API Gateway", COLORS['blue'], 12)

# Services Layer (Teal)
y_svc = y_api + 0.55
add_arrow(slide3, 5, y_svc - 0.05, 5, y_svc + 0.05)
add_label(slide3, "Core Services", LEFT_MARGIN, y_svc + 0.05, 2.5, 0.2, 10)
y_svc_c = y_svc + 0.3
svc_width = (CONTENT_WIDTH - (4 * GAP)) / 5
svcs = ["Member", "Content", "Campaign", "Analytics", "Personalize"]
for i, svc in enumerate(svcs):
    add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(svc_width + GAP), y_svc_c,
                       svc_width, 0.35,
                       svc, COLORS['teal'], 9)

# Data Layer
y_data = y_svc_c + 0.45
add_arrow(slide3, 5, y_data - 0.05, 5, y_data + 0.05)
add_label(slide3, "Data Platform", LEFT_MARGIN, y_data + 0.05, 2, 0.2, 10)
y_data_c = y_data + 0.3
data_width = (CONTENT_WIDTH - (2 * GAP)) / 3
dbs = ["Member DB", "Marketing DB", "CDP/Analytics"]
for i, db in enumerate(dbs):
    add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(data_width + GAP), y_data_c,
                       data_width, 0.35,
                       db, COLORS['green'], 10)

# External Systems (Purple) - Right side
ext_x = LEFT_MARGIN + CONTENT_WIDTH - 1.8
ext_y = 2.2
add_label(slide3, "External\nIntegrations", ext_x, ext_y - 0.3, 1.8, 0.25, 9)
ext_systems = ["Email Service\n(SendGrid)", "CRM\n(Salesforce)", "Social Media"]
for i, sys in enumerate(ext_systems):
    add_shape_with_text(slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
                       ext_x, ext_y + i*0.48,
                       1.8, 0.4,
                       sys, COLORS['purple'], 9)

# Arrows to external systems
for i in range(3):
    add_arrow(slide3, ext_x - 0.05, ext_y + i*0.48 + 0.2,
             LEFT_MARGIN + CONTENT_WIDTH - 1.85, ext_y + i*0.48 + 0.2,
             RGBColor(150, 150, 150))

# Footer
add_label(slide3, "Super Fund MarTech Ecosystem | Future State Architecture",
         LEFT_MARGIN, 5.45, CONTENT_WIDTH, 0.15, 8, False, RGBColor(150, 150, 150))

# Save
output_path = "/Users/touttram/CODE/AAGLOBAL/troy-skills/powerpoint-architecture-skill/superfund-martech-architecture.pptx"
prs.save(output_path)

print(f"✅ Super Fund MarTech diagram created: {output_path}")
print("\nArchitecture Includes:")
print("- Channel Layer: Public Website, Member Portal, iOS App, Android App")
print("- API Gateway: Centralized API management")
print("- MarTech Services: Member Services, Content, Campaign, Marketing Automation, Analytics, Personalization")
print("- Data Platform: Member DB, Marketing DB, Analytics Platform, CDP")
print("- External Integrations: Email Service, CRM, Social Media")
print("\nv1.2 Features Applied:")
print("- ✅ Teal color for core services")
print("- ✅ Consistent 0.5\" margins")
print("- ✅ 9\" content width")
print("- ✅ Full-width layer alignment")
