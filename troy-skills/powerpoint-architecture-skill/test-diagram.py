#!/usr/bin/env python3
"""
Test PowerPoint Architecture Diagram
Demonstrates v1.2 features: Teal color scheme, consistent alignment, full-width layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Define colors (v1.2 color scheme)
COLORS = {
    'blue': RGBColor(0, 120, 212),      # #0078D4 - API Gateway
    'teal': RGBColor(0, 128, 128),      # #008080 - Core Services (NEW in v1.2)
    'orange': RGBColor(255, 140, 0),    # #FF8C00 - UI/Channels
    'green': RGBColor(16, 124, 16),     # #107C10 - Databases
    'purple': RGBColor(92, 45, 145),    # #5C2D91 - External services
    'red': RGBColor(232, 17, 35),       # #E81123 - Security
    'gray': RGBColor(118, 118, 118),    # #767676 - Infrastructure
}

# v1.2 Layout constants - CONSISTENT ALIGNMENT
LEFT_MARGIN = 0.5
RIGHT_MARGIN = 0.5
CONTENT_WIDTH = 9.0  # 10 - 0.5 - 0.5
GAP = 0.1

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color, text_size=14):
    """Add a shape with text and styling"""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    # Fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # Border
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(2)

    # Text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    # Text formatting
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(text_size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)

    return shape

def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(100, 100, 100)):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

title_box = slide1.shapes.add_textbox(
    Inches(LEFT_MARGIN),
    Inches(2),
    Inches(CONTENT_WIDTH),
    Inches(1.5)
)
title_frame = title_box.text_frame
title_frame.text = "Multi-Channel Digital Ecosystem"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

subtitle_box = slide1.shapes.add_textbox(
    Inches(LEFT_MARGIN),
    Inches(3.5),
    Inches(CONTENT_WIDTH),
    Inches(0.5)
)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Future State Architecture - Test Diagram v1.2"
p2 = subtitle_frame.paragraphs[0]
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(100, 100, 100)

# Slide 2: 4-Layer Architecture (Testing v1.2 alignment)
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title = slide2.shapes.add_textbox(
    Inches(LEFT_MARGIN),
    Inches(0.2),
    Inches(CONTENT_WIDTH),
    Inches(0.4)
)
title.text_frame.text = "System Architecture"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True

# Layer 1: Channel Layer (Orange) - 4 components full width
y_layer1 = 0.85
layer1_label = slide2.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(y_layer1), Inches(2), Inches(0.25))
layer1_label.text_frame.text = "Channel Layer"
layer1_label.text_frame.paragraphs[0].font.size = Pt(12)
layer1_label.text_frame.paragraphs[0].font.bold = True

y_layer1_comps = y_layer1 + 0.3
comp_width = (CONTENT_WIDTH - (3 * GAP)) / 4  # 4 equal components

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN, y_layer1_comps, comp_width, 0.5,
                   "Website", COLORS['orange'], 12)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + (comp_width + GAP), y_layer1_comps, comp_width, 0.5,
                   "iOS App", COLORS['orange'], 12)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + 2*(comp_width + GAP), y_layer1_comps, comp_width, 0.5,
                   "Android App", COLORS['orange'], 12)

add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                   LEFT_MARGIN + 3*(comp_width + GAP), y_layer1_comps, comp_width, 0.5,
                   "Admin Portal", COLORS['orange'], 12)

# Arrow down
arrow_y = y_layer1_comps + 0.65
add_arrow(slide2, 5, arrow_y, 5, arrow_y + 0.15)

# Layer 2: API Gateway (Blue) - Full width
y_layer2 = arrow_y + 0.2
layer2_label = slide2.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(y_layer2), Inches(2), Inches(0.25))
layer2_label.text_frame.text = "API Layer"
layer2_label.text_frame.paragraphs[0].font.size = Pt(12)
layer2_label.text_frame.paragraphs[0].font.bold = True

y_layer2_comps = y_layer2 + 0.3
add_shape_with_text(slide2, MSO_SHAPE.HEXAGON,
                   LEFT_MARGIN, y_layer2_comps, CONTENT_WIDTH, 0.5,
                   "API Gateway", COLORS['blue'], 14)

# Arrow down
arrow_y2 = y_layer2_comps + 0.65
add_arrow(slide2, 5, arrow_y2, 5, arrow_y2 + 0.15)

# Layer 3: Core Services (Teal - NEW COLOR in v1.2)
y_layer3 = arrow_y2 + 0.2
layer3_label = slide2.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(y_layer3), Inches(2.5), Inches(0.25))
layer3_label.text_frame.text = "Services Layer (Teal - v1.2)"
layer3_label.text_frame.paragraphs[0].font.size = Pt(12)
layer3_label.text_frame.paragraphs[0].font.bold = True

y_layer3_comps = y_layer3 + 0.3
comp_width3 = (CONTENT_WIDTH - (4 * GAP)) / 5  # 5 equal services

services = ["Auth", "User Mgmt", "Content", "Payment", "Analytics"]
for i, service in enumerate(services):
    add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(comp_width3 + GAP), y_layer3_comps,
                       comp_width3, 0.45,
                       service, COLORS['teal'], 10)

# Arrow down
arrow_y3 = y_layer3_comps + 0.6
add_arrow(slide2, 5, arrow_y3, 5, arrow_y3 + 0.15)

# Layer 4: Data Layer (Green) - 4 databases full width
y_layer4 = arrow_y3 + 0.2
layer4_label = slide2.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(y_layer4), Inches(2), Inches(0.25))
layer4_label.text_frame.text = "Data Layer"
layer4_label.text_frame.paragraphs[0].font.size = Pt(12)
layer4_label.text_frame.paragraphs[0].font.bold = True

y_layer4_comps = y_layer4 + 0.3
comp_width4 = (CONTENT_WIDTH - (3 * GAP)) / 4  # 4 equal databases

databases = ["User DB", "Content DB", "Analytics DB", "Cache"]
for i, db in enumerate(databases):
    add_shape_with_text(slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
                       LEFT_MARGIN + i*(comp_width4 + GAP), y_layer4_comps,
                       comp_width4, 0.4,
                       db, COLORS['green'], 10)

# Footer
footer_y = 5.4
footer = slide2.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(footer_y), Inches(CONTENT_WIDTH), Inches(0.15))
footer.text_frame.text = "v1.2 Features: Teal Services | 0.5\" Margins | 9\" Content Width | Full-Width Alignment"
footer.text_frame.paragraphs[0].font.size = Pt(8)
footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

# Save presentation
output_path = "/Users/touttram/CODE/AAGLOBAL/troy-skills/powerpoint-architecture-skill/test-architecture-v1.2.pptx"
prs.save(output_path)

print(f"✅ Test diagram created: {output_path}")
print("\nv1.2 Features Demonstrated:")
print("- ✅ Teal color (#008080) for Core Services layer")
print("- ✅ Consistent 0.5\" left/right margins")
print("- ✅ 9\" content width across all layers")
print("- ✅ Full-width alignment for all components")
print("- ✅ 4-layer architecture within 5.625\" height limit")
